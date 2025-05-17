#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Präsentationsgenerator aus YAML-Datendatei

Dieses Skript liest Präsentationsdaten aus einer YAML-Datei und erstellt
daraus eine PowerPoint-Präsentation mit Python-PPTX.
"""

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime

def hex_to_rgb(hex_color):
    """Konvertiert eine Hex-Farbe (#RRGGBB) in RGB-Tuple (r, g, b)"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_slide_info_to_presentation(prs, slide_data, styles, settings):
    """Fügt gemeinsame Informationen für Titel und Text zu einer Folie hinzu"""
    slide_layout = prs.slide_layouts[1]  # Title & Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Titel im Apple-Stil
    title_shape = slide.shapes.title
    title_shape.text = slide_data['title']
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(styles['bullet_slide']['title_font_size'])
        paragraph.font.name = settings['font']
    
    # Text im Apple-Stil: klare Hierarchie, ausreichend Weißraum
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, bullet_text in enumerate(slide_data['bullets']):
        if i == 0:
            p = tf.paragraphs[0]
            # Erste Zeile als Kernbotschaft formatieren (fett)
            p.font.bold = styles['bullet_slide']['bullet_main_bold']
            p.font.size = Pt(styles['bullet_slide']['bullet_main_font_size'])
            p.space_after = Pt(styles['bullet_slide']['bullet_spacing_after_main'])
        else:
            p = tf.add_paragraph()
            p.font.size = Pt(styles['bullet_slide']['bullet_sub_font_size'])
            p.space_before = Pt(styles['bullet_slide']['bullet_spacing_before_sub'])
        
        p.text = bullet_text
        p.font.name = settings['font']

    return slide

def add_title_slide(prs, slide_data, styles, settings):
    """Fügt eine Titelfolie hinzu"""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Großer, klarer Titel im Apple-Stil
    title_shape = slide.shapes.title
    title_shape.text = slide_data['title']
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(styles['title_slide']['title_font_size'])
        paragraph.font.bold = styles['title_slide']['title_font_bold']
        paragraph.font.name = settings['font']
    
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = slide_data['subtitle']
    for paragraph in subtitle_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(styles['title_slide']['subtitle_font_size'])
        paragraph.font.name = settings['font']

def add_bullet_slide(prs, slide_data, styles, settings):
    """Fügt eine Aufzählungsfolie hinzu"""
    slide = add_slide_info_to_presentation(prs, slide_data, styles, settings)
    
    if 'placeholder' in slide_data:
        # Platzhalter für Bilder/Diagramme
        add_placeholder(slide, slide_data['placeholder'], styles, settings)

def add_result_slide(prs, slide_data, styles, settings):
    """Fügt eine Ergebnisfolie mit KPI-Badge hinzu"""
    slide = add_slide_info_to_presentation(prs, slide_data, styles, settings)
    
    # KPI-Badge im Apple-Stil
    kpi_style = styles['result_slide']['kpi_badge']
    left, top, width, height = Inches(0.5), Inches(5.5), Inches(kpi_style['width']), Inches(kpi_style['height'])
    kpi_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    kpi_shape.fill.solid()
    
    # Farbe aus Hex-Wert umwandeln wenn nötig
    if isinstance(kpi_style['color'], str) and kpi_style['color'].startswith('#'):
        r, g, b = hex_to_rgb(kpi_style['color'])
        kpi_shape.fill.fore_color.rgb = RGBColor(r, g, b)
    else:
        kpi_shape.fill.fore_color.rgb = RGBColor(40, 167, 69)  # Fallback: Sattes Grün
        
    kpi_shape.line.width = Pt(0)  # Keine Umrandung für modernen Look
    
    # Text mit ausreichend Weißraum
    text_frame = kpi_shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    text_frame.margin_top = Inches(0.1)
    text_frame.vertical_anchor = 1  # Mittig
    
    p = text_frame.add_paragraph()
    p.text = slide_data['kpi_text']
    
    # Textfarbe aus Hex-Wert umwandeln wenn nötig
    if isinstance(kpi_style['font_color'], str) and kpi_style['font_color'].startswith('#'):
        r, g, b = hex_to_rgb(kpi_style['font_color'])
        p.font.color.rgb = RGBColor(r, g, b)
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)  # Fallback: Weiß
        
    p.font.size = Pt(kpi_style['font_size'])
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = settings['font']

    if 'placeholder' in slide_data:
        # Platzhalter für Bilder/Diagramme
        add_placeholder(slide, slide_data['placeholder'], styles, settings)

def add_placeholder(slide, placeholder_data, styles, settings):
    """Fügt einen Platzhalter für Bilder/Diagramme hinzu"""
    ph_style = styles['placeholder']
    left, top, width, height = Inches(8.8), Inches(1.8), Inches(ph_style['width']), Inches(ph_style['height'])
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    ph.fill.solid()
    
    # Farben aus Hex-Werten umwandeln wenn nötig
    if isinstance(ph_style['color'], str) and ph_style['color'].startswith('#'):
        r, g, b = hex_to_rgb(ph_style['color'])
        ph.fill.fore_color.rgb = RGBColor(r, g, b)
    else:
        ph.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Fallback: Sehr helles Grau
        
    if isinstance(ph_style['border_color'], str) and ph_style['border_color'].startswith('#'):
        r, g, b = hex_to_rgb(ph_style['border_color'])
        ph.line.color.rgb = RGBColor(r, g, b)
    else:
        ph.line.color.rgb = RGBColor(220, 220, 220)  # Fallback: Subtiler Rahmen
        
    ph.shadow.inherit = False  # Kein Schatten (klarer Apple-Look)
    
    # Den Platzhalter in den Hintergrund verschieben
    ph.zorder = 0
    
    # Placeholder-Text
    text_frame = ph.text_frame
    text_frame.word_wrap = True
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.vertical_anchor = 1  # Mittig
    
    # Haupttext
    p = text_frame.add_paragraph()
    p.text = placeholder_data.get('text', "Platzhalter")
    p.font.size = Pt(ph_style['font_size'])
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = settings['font']
    
    # Beschreibungstext falls vorhanden
    if 'note' in placeholder_data:
        p = text_frame.add_paragraph()
        p.text = placeholder_data['note']
        p.font.size = Pt(ph_style['note_font_size'])
        p.alignment = PP_ALIGN.CENTER
        p.font.italic = True
        p.font.name = settings['font']
        
        # Textfarbe aus Hex-Wert umwandeln wenn nötig
        if isinstance(ph_style['note_font_color'], str) and ph_style['note_font_color'].startswith('#'):
            r, g, b = hex_to_rgb(ph_style['note_font_color'])
            p.font.color.rgb = RGBColor(r, g, b)
        else:
            p.font.color.rgb = RGBColor(100, 100, 100)  # Fallback: Dunkelgrau

def load_presentation_data(yaml_file):
    """Lädt die Präsentationsdaten aus einer YAML-Datei"""
    with open(yaml_file, 'r', encoding='utf-8') as file:
        data = yaml.safe_load(file)
    return data['presentation']

def create_presentation(yaml_file, output_path=None):
    """Erstellt eine Präsentation aus den Daten in einer YAML-Datei"""
    # Daten laden
    presentation_data = load_presentation_data(yaml_file)
    settings = presentation_data['settings']
    styles = presentation_data['styles']
    
    # Neue Präsentation erstellen
    prs = Presentation()
    
    # Globale Einstellungen
    prs.slide_width = Inches(settings['width'])
    prs.slide_height = Inches(settings['height'])
    
    # Folien erstellen
    for slide_data in presentation_data['slides']:
        if slide_data['type'] == 'title':
            add_title_slide(prs, slide_data, styles, settings)
        elif slide_data['type'] == 'bullet':
            add_bullet_slide(prs, slide_data, styles, settings)
        elif slide_data['type'] == 'result':
            add_result_slide(prs, slide_data, styles, settings)
    
    # Wenn kein Ausgabepfad angegeben, auf dem Desktop speichern
    if not output_path:
        desktop_path = os.path.expanduser("~/Desktop")
        today = datetime.now().strftime('%Y-%m-%d')
        filename = f"Freelance_Success_Stories_Michael_Boiman_{today}.pptx"
        output_path = os.path.join(desktop_path, filename)
    
    # Präsentation speichern
    prs.save(output_path)
    
    return output_path

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Erstellt eine PowerPoint-Präsentation aus einer YAML-Datendatei')
    parser.add_argument('-i', '--input', default='presentation_data.yaml',
                        help='Pfad zur YAML-Datendatei (Standard: presentation_data.yaml)')
    parser.add_argument('-o', '--output', 
                        help='Pfad zur Ausgabedatei (Standard: Desktop/Freelance_Success_Stories_Michael_Boiman_YYYY-MM-DD.pptx)')
    
    args = parser.parse_args()
    
    try:
        # Präsentationsdaten für Font-Information laden
        presentation_data = load_presentation_data(args.input)
        font_name = presentation_data['settings']['font']
        
        # Präsentation erzeugen
        output_path = create_presentation(args.input, args.output)
        print(f"Präsentation wurde erfolgreich gespeichert unter: {output_path}")
        print(f"HINWEIS: Die Präsentation verwendet die Schriftart '{font_name}'. "
              f"Falls diese auf Ihrem System nicht verfügbar ist, wird PowerPoint/Keynote automatisch eine Alternative verwenden.")
    except Exception as e:
        print(f"Fehler beim Erstellen der Präsentation: {e}")
