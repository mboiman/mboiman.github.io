from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime

# ---------- Helper functions ----------
def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Großer, klarer Titel im Apple-Stil
    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.name = "SF Pro Display"  # Apple-ähnliche Schriftart
    
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    for paragraph in subtitle_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.font.name = "SF Pro Display"

def add_bullet_slide(prs, title, bullets, placeholder=False, placeholder_text="Screenshot/Diagramm hier einfügen", placeholder_note=None):
    slide_layout = prs.slide_layouts[1]  # Title & Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Titel im Apple-Stil
    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.name = "SF Pro Display"
    
    # Text im Apple-Stil: klare Hierarchie, ausreichend Weißraum
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            # Erste Zeile als Kernbotschaft formatieren (fett)
            p.font.bold = True
            p.font.size = Pt(24)  # Größere Schrift für Hauptaussage
            p.space_after = Pt(14)  # Mehr Abstand nach der Hauptaussage
        else:
            p = tf.add_paragraph()
            p.font.size = Pt(20)  # Gut lesbare Größe für Bullet-Points
            p.space_before = Pt(10)  # Abstand vor jedem Bullet-Point
        
        p.text = b
        p.font.name = "SF Pro Display"  # Apple-ähnliche Schriftart

    if placeholder:
        # Position für besseres visuelles Gleichgewicht anpassen
        left, top, width, height = Inches(5.3), Inches(1.8), Inches(4), Inches(4)
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Sehr helles Grau (Apple-Stil)
        ph.line.color.rgb = RGBColor(220, 220, 220)  # Subtiler Rahmen
        ph.shadow.inherit = False  # Kein Schatten (klarer Apple-Look)
        
        # Placeholder-Text im Apple-Stil
        text_frame = ph.text_frame
        text_frame.word_wrap = True
        text_frame.margin_bottom = Inches(0.1)
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.vertical_anchor = 1  # Mittig
        
        # Haupttext
        p = text_frame.add_paragraph()
        p.text = placeholder_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "SF Pro Display"
        
        # Beschreibungstext falls vorhanden
        if placeholder_note:
            p = text_frame.add_paragraph()
            p.text = placeholder_note
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            p.font.italic = True
            p.font.name = "SF Pro Display"
            p.font.color.rgb = RGBColor(100, 100, 100)  # Dunkelgrau für sekundären Text

# Neue Funktion für Erfolgsfolien mit KPI-Badge
def add_result_slide(prs, title, bullets, kpi_text, placeholder=False, placeholder_text="Screenshot/Diagramm hier einfügen", placeholder_note=None):
    slide_layout = prs.slide_layouts[1]  # Title & Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Titel im Apple-Stil
    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.name = "SF Pro Display"
    
    # Text im Apple-Stil
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            # Erste Zeile als Kernbotschaft formatieren
            p.font.bold = True
            p.font.size = Pt(24)  # Größere Schrift für Hauptaussage
            p.space_after = Pt(14)  # Mehr Abstand nach der Hauptaussage
        else:
            p = tf.add_paragraph()
            p.font.size = Pt(20)  # Gut lesbare Größe für Bullet-Points
            p.space_before = Pt(10)  # Abstand vor jedem Bullet-Point
        
        p.text = b
        p.font.name = "SF Pro Display"
    
    # KPI-Badge im Apple-Stil
    # Modernes, abgerundetes Badge mit mehr Weißraum
    left, top, width, height = Inches(0.5), Inches(5.5), Inches(4.5), Inches(0.7)
    kpi_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    kpi_shape.fill.solid()
    kpi_shape.fill.fore_color.rgb = RGBColor(40, 167, 69)  # Sattes Grün (ähnlich Apple Success)
    kpi_shape.line.width = Pt(0)  # Keine Umrandung für modernen Look
    
    # Text mit ausreichend Weißraum
    text_frame = kpi_shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    text_frame.margin_top = Inches(0.1)
    text_frame.vertical_anchor = 1  # Mittig
    
    p = text_frame.add_paragraph()
    p.text = kpi_text
    p.font.color.rgb = RGBColor(255, 255, 255)  # Weiß
    p.font.size = Pt(24)  # Gut sichtbare Größe
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "SF Pro Display"

    if placeholder:
        # Position für besseres visuelles Gleichgewicht
        left, top, width, height = Inches(5.3), Inches(1.8), Inches(4), Inches(4)
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Sehr helles Grau (Apple-Stil)
        ph.line.color.rgb = RGBColor(220, 220, 220)  # Subtiler Rahmen
        ph.shadow.inherit = False  # Kein Schatten
        
        # Placeholder-Text
        text_frame = ph.text_frame
        text_frame.word_wrap = True
        text_frame.margin_bottom = Inches(0.1)
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.vertical_anchor = 1  # Mittig
        
        # Haupttext
        p = text_frame.add_paragraph()
        p.text = placeholder_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "SF Pro Display"
        
        # Beschreibungstext falls vorhanden
        if placeholder_note:
            p = text_frame.add_paragraph()
            p.text = placeholder_note
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            p.font.italic = True
            p.font.name = "SF Pro Display"
            p.font.color.rgb = RGBColor(100, 100, 100)  # Dunkelgrau

# ---------- Build presentation ----------
prs = Presentation()

# Setze globale Präsentationseinstellungen für klares Design
prs.slide_width = Inches(13.33)  # 16:9 Seitenverhältnis (moderner)
prs.slide_height = Inches(7.5)

# 1 – Titel
add_title_slide(
    prs,
    "Freelance Success Stories – KI‑gestützte Qualität & Automatisierung",
    "Michael Boiman · Test‑ & Automatisierungs‑Experte"
)

# 2 – Agenda
add_bullet_slide(
    prs,
    "Agenda & Lesehinweis",
    ["So lesen Sie diese Präsentation: Jedes Projekt zeigt Problem → Lösung → messbarer Erfolg (mit grünem KPI-Badge)",
     "1. Ausgangssituation & relevante Herausforderungen",
     "2. Vier erfolgreiche Projekte mit nachweisbarem Impact",
     "3. Übergreifender Business Impact & Lessons Learned",
     "4. Kontaktmöglichkeiten für weitere Details oder Live-Demo"]
)

# 3 – Ausgangssituation
add_bullet_slide(
    prs,
    "Ausgangssituation",
    ["In allen Projekten bestand ein konkreter Business-Bedarf nach mehr Effizienz und Transparenz:",
     "• Fragmentierte Tests & hoher manueller Aufwand bei Go/No-Go-Entscheidungen",
     "• Legacy‑Mainframe-Preiskalkulation ohne Nachweisbarkeit",
     "• Wachsende E‑Mail‑Flut mit steigendem SLA‑Risiko",
     "• Ziel: Geschwindigkeit, Transparenz und messbare Qualität"],
    placeholder=True,
    placeholder_text="Pain‑Points‑Grafik",
    placeholder_note="Hier: Visualisierung der Kernprobleme mit Icon für jede Herausforderung"
)

# 4 – Quality Dashboard – Herausforderung
add_bullet_slide(
    prs,
    "PROJEKT 1: Quality Dashboard – Herausforderung",
    ["Ohne Gesamtüberblick war die Testqualität nicht messbar und Go/No-Go-Entscheidungen wurden verzögert:",
     "• >10 Test‑Tools, 5 Umgebungen, kein Gesamtstatus",
     "• Manuelle Rückfragen kosteten Zeit (mehrere Stunden pro Release)",
     "• Management hatte keine Entscheidungsgrundlage für Releases"],
    placeholder=True,
    placeholder_text="Vorher-Situation",
    placeholder_note="Hier: Chaotisches Dashboard/Whiteboard mit vielen verschiedenen Reports"
)

# 5 – Quality Dashboard – Lösung
add_bullet_slide(
    prs,
    "PROJEKT 1: Quality Dashboard – Meine Lösung",
    ["Ich entwickelte eine vollautomatische End-to-End-Lösung für Echtzeit-Qualitätsübersicht:",
     "• Pipeline: Reports aus allen Tools → Parser → zentrale DB → Dashboard/PDF",
     "• Echtzeit‑Filter nach Umgebung, Projekt, Applikation",
     "• Automatische Jira-Ticket-Erstellung bei kritischen Fehlern",
     "• Customizable Widgets für verschiedene Stakeholder"],
    placeholder=True,
    placeholder_text="Flow‑Diagramm",
    placeholder_note="Hier: Architekturdiagramm der Pipeline-Lösung mit Tools und Datenfluss"
)

# 6 – Quality Dashboard – Ergebnis
add_result_slide(
    prs,
    "PROJEKT 1: Quality Dashboard – Ergebnis & Impact",
    ["Die Lösung revolutionierte den Release-Prozess und steigerte die Entwicklungsgeschwindigkeit:", 
     "• Live‑Status in unter 5 Sekunden für alle Teams verfügbar",
     "• Trendanalysen über Historie ermöglichen Qualitätsprognosen",
     "• PDF-Reports auf Knopfdruck für Management-Meetings"],
    "IMPACT: -80% Rückfragen | Go/No-Go-Entscheidungen in Minuten statt Stunden",
    placeholder=True,
    placeholder_text="Dashboard‑Screenshot",
    placeholder_note="Hier: Screenshot des Live-Dashboards mit Qualitätsampeln und Trends"
)

# 7 – Dual‑Run Preiskalkulation – Problem
add_bullet_slide(
    prs,
    "PROJEKT 2: Dual‑Run Preiskalkulation – Problem",
    ["Die Ablösung eines kritischen Legacy-Mainframe-Systems stellte besondere Herausforderungen:",
     "• >1 Mio Kombinationsmöglichkeiten in der Preisberechnung",
     "• Kein vollständiges Testset möglich (zu komplex für manuelle Tests)",
     "• Hoher Druck: Fachbereiche brauchten Beweise für korrekte Funktion"],
    placeholder=True,
    placeholder_text="Parameter‑Diagramm",
    placeholder_note="Hier: Visualisierung der Parametervielfalt und Komplexität"
)

# 8 – Dual‑Run Preiskalkulation – Teststrategie
add_bullet_slide(
    prs,
    "PROJEKT 2: Dual‑Run Preiskalkulation – Meine Lösung",
    ["Ich entwickelte einen innovativen 'Dual-Run'-Ansatz mit kontinuierlichem Vergleich:",
     "• Zufallsgenerator erstellt und testet komplexe Szenarien 24/7",
     "• Parallele Ausführung in Alt- und Neusystem mit pixelgenauem Vergleich",
     "• Automatischer Ergebnisvergleich mit Drill-Down-Funktion im Dashboard",
     "• Self-Service für Fachanwender: Fehleranalyse ohne Entwicklersupport"],
    placeholder=True,
    placeholder_text="Sequence‑Diagramm",
    placeholder_note="Hier: Visualisierung des Dual-Run-Prozesses mit Vergleichslogik"
)

# 9 – Dual‑Run Preiskalkulation – Impact
add_result_slide(
    prs,
    "PROJEKT 2: Dual‑Run Preiskalkulation – Ergebnis",
    ["Die Lösung ermöglichte die sichere Migration kritischer Business-Funktionen:",
     "• Automatisierte Tests deckten versteckte Fehler in beiden Systemen auf",
     "• Fachbereich kann selbständig jede Abweichung bis zur Einzelposition nachvollziehen",
     "• Management erhält verlässliche Migrationsmetriken für Go-Live-Entscheidung"],
    "IMPACT: Qualitätskennzahl von 30% auf 96% erhöht | Direkte Fehlerlokalisierung",
    placeholder=True,
    placeholder_text="Progress‑Chart",
    placeholder_note="Hier: Fortschrittskurve mit steigender Übereinstimmungsrate 30% → 96%"
)

# 10 – KI‑Testautomatisierung – Approach
add_bullet_slide(
    prs,
    "PROJEKT 3: KI‑Testautomatisierung – Herausforderung & Lösungsansatz",
    ["Testautomatisierung war zu zeitaufwändig und inkonsistent - mein KI-Ansatz löste das:",
     "• Problem: Langsame Testerstellung und unterschiedliche Code-Qualität",
     "• Meine Lösung: Guideline-basierter KI-Agent im IDE für konsistenten Code",
     "• Erreicht: 80% des Test-Codes automatisch generiert",
     "• Innovation: Prompt-basierte Testgenerierung direkt aus Jira-Stories (60%)"],
    placeholder=True,
    placeholder_text="IDE‑Screenshot",
    placeholder_note="Hier: Screenshot des KI-Agenten bei der Code-Generierung im IDE"
)

# 11 – KI‑Testautomatisierung – Nutzen
add_result_slide(
    prs,
    "PROJEKT 3: KI‑Testautomatisierung – Ergebnis & Business Value",
    ["Der KI-Ansatz transformierte die Testentwicklung im gesamten Unternehmen:",
     "• Konsistente Code-Qualität über alle Teams hinweg",
     "• Verständliche Fachreports mit Screenshots für Nicht-Techniker",
     "• Self-Service-Reports für Fachbereiche ohne IT-Support"],
    "IMPACT: 50% weniger Entwicklungszeit pro Test-Case | Review-Aufwand -40%",
    placeholder=True,
    placeholder_text="Testreport‑Ausschnitt",
    placeholder_note="Hier: Beispiel eines automatisch generierten Reports mit Screenshots"
)

# 12 – KI‑E‑Mail‑Prozess – Workflow
add_bullet_slide(
    prs,
    "PROJEKT 4: KI‑E‑Mail‑Prozess – Herausforderung & Lösung",
    ["Hohe manuelle E-Mail-Bearbeitungszeit verursachte SLA-Risiken - meine KI-Lösung:",
     "• Problem: Große E-Mail-Mengen mit komplexem Routing und langer Bearbeitungszeit",
     "• Meine Lösung: End-to-End KI-Pipeline mit automatischer Klassifizierung",
     "• Prozess: Eingang → KI-Klassifikation → SAP/Salesforce-Integration → Teams-Notification",
     "• Audit-Sicherheit: Revisionssichere Logs & Echtzeit-Dashboard"],
    placeholder=True,
    placeholder_text="Swimlane‑Diagramm",
    placeholder_note="Hier: Workflow-Diagramm des E-Mail-Verarbeitungsprozesses"
)

# 13 – KI‑E‑Mail‑Prozess – Zahlen & Nutzen
add_result_slide(
    prs,
    "PROJEKT 4: KI‑E‑Mail‑Prozess – Ergebnis & Skalierung",
    ["Die Lösung eliminierte manuelle Verarbeitung und skalierte problemlos:",
     "• Vollautomatische E-Mail- und Anhang-Verarbeitung ohne menschlichen Eingriff",
     "• Echtzeit-Dashboard für SLA-Tracking und Lastverteilung",
     "• Problemlose Skalierung auf zweites Unternehmen ohne Mehraufwand"],
    "IMPACT: Bearbeitungszeit von Minuten auf Sekunden | 100% Audit-Trail",
    placeholder=True,
    placeholder_text="Heatmap‑Dashboard",
    placeholder_note="Hier: Dashboard mit Heatmap nach Volumen/Stadt und SLA-Metriken"
)

# 14 – Business Impact (Gesamt)
add_bullet_slide(
    prs,
    "Business Impact (Gesamt)",
    ["Die kombinierten Lösungen lieferten signifikanten, messbaren Geschäftswert:",
     "• >120.000 € jährliche Einsparung durch Automatisierung & Effizienzsteigerung",
     "• Beschleunigte Release‑Entscheidungen: Von 2 Tagen auf 1 Stunde",
     "• Entlastung der Fachbereiche durch Self-Service & Transparenz",
     "• Qualitätssteigerung: Messbar bessere Software bei schnelleren Releases"],
    placeholder=True,
    placeholder_text="Balkendiagramm",
    placeholder_note="Hier: Visualisierung des Business-Impact in Zeit, Qualität & Kosten"
)

# 15 – Lessons Learned
add_bullet_slide(
    prs,
    "Lessons Learned & Best Practices",
    ["Aus allen Projekten gewonnene Erkenntnisse für zukünftige Vorhaben:",
     "• Dashboards iterativ mit Stakeholdern entwickeln, nicht im Big-Bang",
     "• KI braucht klare Guidelines für konsistente Ergebnisse",
     "• Fachbereiche früh einbinden und Self-Service ermöglichen",
     "• Kommunikation wichtiger als perfekte technische Lösung"],
    placeholder=True,
    placeholder_text="Idea‑Icon",
    placeholder_note="Hier: Visualisierung der Key Learnings als Icons/Symbole"
)

# 16 – Kontakt & Demos
add_bullet_slide(
    prs,
    "Kontakt & Live-Demos",
    ["Ich freue mich über die Gelegenheit, diese Lösungen live zu demonstrieren:",
     "• Live‑Demo (Textanalyse): demo‑app‑text‑analyze‑wdrcmwwnk7r9gsnaftxczk.streamlit.app",
     "• Elevator‑Pitch: demo‑app‑text‑analyze‑ehjk3syc4magebt7yvcydt2.streamlit.app",
     "• LinkedIn: linkedin.com/in/michael-boiman-ab709912",
     "• E‑Mail: mboiman@gmail.com",
     "• Website: boiman-kupermann.com",
     "\nLassen Sie uns gemeinsam Ihre Qualitäts- und Automatisierungsziele erreichen!"],
    placeholder=True,
    placeholder_text="QR‑Codes",
    placeholder_note="Hier: QR-Codes zu LinkedIn-Profil und Live-Demos"
)

# ---------- Save presentation ----------
# Erstelle Desktop-Ordner, falls er nicht existiert
desktop_path = os.path.expanduser("~/Desktop")
filename = f"Freelance_Success_Stories_Michael_Boiman_{datetime.now().strftime('%Y-%m-%d')}.pptx"
file_path = os.path.join(desktop_path, filename)

prs.save(file_path)
print(f"Präsentation wurde erfolgreich gespeichert unter: {file_path}")
print(f"HINWEIS: Die Präsentation verwendet die Schriftart 'SF Pro Display'. Falls diese auf Ihrem System nicht verfügbar ist, wird PowerPoint/Keynote automatisch eine Alternative verwenden.")
